"""
draw.io XML → PowerPoint 转换器

将 draw.io 的 mxGraph XML 转为原生可编辑 PPT 形状。
支持 group 嵌套（相对坐标解析）和 HTML 格式文本。

Usage:
    python drwaio2ppt.py
    python drwaio2ppt.py -i input.drawio
    python drwaio2ppt.py -i input.xml -o output.pptx
"""

from __future__ import annotations

import argparse
import re
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ── 常量 ────────────────────────────────────────────────────────────────
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.25
FONT_BOOST = 1.20


# ═════════════════════════════════════════════════════════════════════════
# 文本清洗
# ═════════════════════════════════════════════════════════════════════════

_HTML_TAG = re.compile(r"<[^>]*>")
_HTML_ENTITIES = {"&#xa;": "\n", "&lt;": "<", "&gt;": ">", "&amp;": "&",
                   "&nbsp;": " ", "&quot;": '"'}


def _clean_value(raw: str) -> str:
    """去掉 HTML 标签，保留换行和纯文本。"""
    s = raw
    for ent, ch in _HTML_ENTITIES.items():
        s = s.replace(ent, ch)
    s = _HTML_TAG.sub("", s)  # 去标签
    return s.strip()


# ═════════════════════════════════════════════════════════════════════════
# 样式解析
# ═════════════════════════════════════════════════════════════════════════


def _parse_style(style_str: str) -> dict[str, str]:
    props: dict[str, str] = {}
    if not style_str:
        return props
    for part in style_str.split(";"):
        part = part.strip()
        if "=" in part:
            k, v = part.split("=", 1)
            props[k.strip()] = v.strip()
        elif part:
            props[part] = "1"
    return props


def _hex_to_rgb(hex_color: str) -> RGBColor:
    c = (hex_color or "").lstrip("#")
    if not c or c == "none":
        return RGBColor(0xFF, 0xFF, 0xFF)
    if len(c) == 3:
        c = "".join(ch * 2 for ch in c)
    if len(c) != 6:
        return RGBColor(0xCC, 0xCC, 0xCC)
    return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))


# ═════════════════════════════════════════════════════════════════════════
# XML 解析 — 支持 group 嵌套相对坐标
# ═════════════════════════════════════════════════════════════════════════


def _classify(style: dict, w: float, h: float) -> str:
    shape = style.get("shape", "")
    dashed = style.get("dashed") == "1"
    fill = style.get("fillColor", "")
    stroke = style.get("strokeColor", "")

    if shape.startswith("cylinder"):
        return "cylinder"
    if shape == "ellipse":
        return "ellipse"
    if shape == "group":
        return "group"
    if stroke == "none" and fill in ("none", "", None):
        return "text"
    if dashed and w > 200:
        return "dashed_box"
    if w > 500 and h > 60:
        return "bg"
    if w < 105 and h > 60:
        return "label"
    return "box"


def parse_drawio(path: str) -> tuple[list[dict], dict, dict]:
    """解析 draw.io 文件。

    关键处理：
    1. 递归解析 parent 引用，将 group 子元素的相对坐标转为绝对坐标
    2. 清洗 HTML 标签保留纯文本
    3. 跳过纯 group 容器（无实际内容）
    """
    tree = ET.parse(path)
    cells = tree.findall(".//mxCell")
    if not cells:
        raise ValueError("未找到 mxCell，不是 draw.io 文件")

    # 第一遍：建立 id → { 原始 geo, parent_id, style, value, waypoints, 子节点列表 }
    nodes: dict[str, dict] = {}
    for cell in cells:
        cid = cell.get("id", "")
        if cid in ("0", "1"):
            continue
        geo = cell.find("mxGeometry")
        if geo is None:
            continue
        style = _parse_style(cell.get("style", ""))
        value = _clean_value(cell.get("value") or "")

        # 提取 waypoints（正交走线拐点）
        waypoints: list[tuple[float, float]] = []
        arr = geo.find("Array")
        if arr is not None:
            for pt in arr.findall("mxPoint"):
                waypoints.append((float(pt.get("x", 0)), float(pt.get("y", 0))))

        # sourcePoint / targetPoint（绝对坐标端点，无 source/target 引用时使用）
        sp = geo.find("sourcePoint")
        tp = geo.find("targetPoint")
        spt = (float(sp.get("x", 0)), float(sp.get("y", 0))) if sp is not None else None
        tpt = (float(tp.get("x", 0)), float(tp.get("y", 0))) if tp is not None else None

        nodes[cid] = {
            "geo": {
                "x": float(geo.get("x", 0)),
                "y": float(geo.get("y", 0)),
                "w": float(geo.get("width", 0)),
                "h": float(geo.get("height", 0)),
            },
            "parent": cell.get("parent", "1"),
            "style": style,
            "value": value,
            "edge": cell.get("edge") == "1",
            "source": cell.get("source"),
            "target": cell.get("target"),
            "source_pt": spt,
            "target_pt": tpt,
            "waypoints": waypoints,
            "children": [],
        }

    # 建立父子关系
    for cid, nd in nodes.items():
        pid = nd["parent"]
        if pid in nodes:
            nodes[pid]["children"].append(cid)

    # 第二遍：递归计算绝对坐标
    def resolve_abs(cid: str, acc_x: float = 0, acc_y: float = 0) -> None:
        nd = nodes[cid]
        g = nd["geo"]
        g["x"] += acc_x
        g["y"] += acc_y
        for child_id in nd["children"]:
            resolve_abs(child_id, g["x"], g["y"])

    for cid in nodes:
        nd = nodes[cid]
        if nd["parent"] in ("1", "0"):  # 顶级元素
            resolve_abs(cid, 0, 0)

    # 第三遍：构建形状列表（跳过 group 容器）
    shapes: list[dict] = []
    id_geo: dict[str, dict] = {}
    min_x, min_y = float("inf"), float("inf")
    max_x, max_y = float("-inf"), float("-inf")

    for cid, nd in nodes.items():
        g = nd["geo"]
        id_geo[cid] = {"x": g["x"], "y": g["y"], "w": g["w"], "h": g["h"]}

        if nd["edge"]:
            shapes.append({
                "type": "edge", "value": nd["value"],
                "x": g["x"], "y": g["y"], "w": g["w"], "h": g["h"],
                "style": nd["style"],
                "source": nd["source"], "target": nd["target"],
                "source_pt": nd["source_pt"], "target_pt": nd["target_pt"],
                "waypoints": nd["waypoints"],
            })
            continue

        s_type = _classify(nd["style"], g["w"], g["h"])
        if s_type == "group":
            continue  # 不绘制 group 容器本身

        if g["w"] > 0 and g["h"] > 0:
            min_x = min(min_x, g["x"]); min_y = min(min_y, g["y"])
            max_x = max(max_x, g["x"] + g["w"]); max_y = max(max_y, g["y"] + g["h"])

        shapes.append({
            "type": s_type,
            "value": nd["value"],
            "x": g["x"], "y": g["y"], "w": g["w"], "h": g["h"],
            "style": nd["style"], "source": None, "target": None,
        })

    bounds = {
        "x": min_x if min_x != float("inf") else 0,
        "y": min_y if min_y != float("inf") else 0,
        "w": max_x - min_x if max_x != float("-inf") else 800,
        "h": max_y - min_y if max_y != float("-inf") else 600,
    }
    return shapes, bounds, id_geo


# ═════════════════════════════════════════════════════════════════════════
# PPT 转换器
# ═════════════════════════════════════════════════════════════════════════


class Converter:
    def __init__(self):
        self.scale = 1.0
        self.off_x = 0.0
        self.off_y = 0.0

    def _init(self, bounds: dict) -> None:
        cw_in = bounds["w"] / 72.0
        ch_in = bounds["h"] / 72.0
        avail_w = SLIDE_W - 2 * MARGIN
        avail_h = SLIDE_H - 2 * MARGIN
        self.scale = min(avail_w / cw_in, avail_h / ch_in) if cw_in > 0 else 1.0
        scaled_w = cw_in * self.scale
        scaled_h = ch_in * self.scale
        self.off_x = ((SLIDE_W - scaled_w) / 2) - (bounds["x"] / 72.0 * self.scale)
        self.off_y = ((SLIDE_H - scaled_h) / 2) - (bounds["y"] / 72.0 * self.scale)

    def _x(self, du: float) -> int:
        return int((du / 72.0 * self.scale + self.off_x) * 914400)

    def _y(self, du: float) -> int:
        return int((du / 72.0 * self.scale + self.off_y) * 914400)

    def _w(self, du: float) -> int:
        return int(du / 72.0 * self.scale * 914400)

    def _h(self, du: float) -> int:
        return int(du / 72.0 * self.scale * 914400)

    def _font_pt(self, drawio_pt: float) -> float:
        return round(drawio_pt * self.scale * FONT_BOOST, 1)

    # ── 文字设置 ──────────────────────────────────────────────────────

    def _tf_cfg(self, shape, text: str, *, font_size: float = 10,
                bold: bool = False, color: RGBColor | None = None,
                halign: PP_ALIGN = PP_ALIGN.CENTER,
                vanchor: MSO_ANCHOR = MSO_ANCHOR.MIDDLE) -> None:
        if not text:
            return
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        tf.margin_top = Pt(1)
        tf.margin_bottom = Pt(1)
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = halign
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(font_size)
            r.font.bold = bold
            if color:
                r.font.color.rgb = color

    # ── 逐类型绘制 ────────────────────────────────────────────────────

    def _bg(self, slide, s):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self._x(s["x"]), self._y(s["y"]),
            self._w(s["w"]), self._h(s["h"]),
        )
        fill = _hex_to_rgb(s["style"].get("fillColor", "#eee"))
        stroke = _hex_to_rgb(s["style"].get("strokeColor", "#999"))
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = stroke; shape.line.width = Pt(1)

    def _label(self, slide, s):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self._x(s["x"]), self._y(s["y"]),
            self._w(s["w"]), self._h(s["h"]),
        )
        fill = _hex_to_rgb(s["style"].get("fillColor", "#ddd"))
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = RGBColor(0x36, 0x39, 0x3D)
        shape.line.width = Pt(0.5)
        self._tf_cfg(shape, s["value"], font_size=self._font_pt(12),
                     bold=True, color=RGBColor(0x36, 0x39, 0x3D))

    def _dashed_box(self, slide, s):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self._x(s["x"]), self._y(s["y"]),
            self._w(s["w"]), self._h(s["h"]),
        )
        stroke = _hex_to_rgb(s["style"].get("strokeColor", "#999"))
        fill_hex = s["style"].get("fillColor", "none")
        if fill_hex == "none":
            shape.fill.background()
        else:
            shape.fill.solid(); shape.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
        shape.line.color.rgb = stroke; shape.line.width = Pt(0.8)
        spPr = shape._element.find(qn("a:spPr"))
        if spPr is not None:
            ln = spPr.find(qn("a:ln"))
            if ln is None:
                ln = ET.SubElement(spPr, qn("a:ln"))
            ET.SubElement(ln, qn("a:prstDash")).set("val", "dash")
        va = s["style"].get("verticalAlign", "middle")
        vanchor = MSO_ANCHOR.TOP if va == "top" else MSO_ANCHOR.MIDDLE
        self._tf_cfg(shape, s["value"], font_size=self._font_pt(8),
                     bold=True, color=stroke, halign=PP_ALIGN.LEFT,
                     vanchor=vanchor)

    def _text(self, slide, s):
        if not s["value"]:
            return
        shape = slide.shapes.add_textbox(
            self._x(s["x"]), self._y(s["y"]),
            self._w(s["w"]), self._h(s["h"]),
        )
        fs = int(s["style"].get("fontSize", "10"))
        bold = s["style"].get("fontStyle") == "1"
        ha = s["style"].get("align", "left")
        halign = PP_ALIGN.LEFT if ha == "left" else PP_ALIGN.CENTER
        self._tf_cfg(shape, s["value"], font_size=self._font_pt(fs),
                     bold=bold, color=RGBColor(0x36, 0x39, 0x3D),
                     halign=halign)

    def _cylinder(self, slide, s):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.FLOWCHART_PUNCHED_TAPE,
            self._x(s["x"]), self._y(s["y"]),
            self._w(s["w"]), self._h(s["h"]),
        )
        fill = _hex_to_rgb(s["style"].get("fillColor", "#fff"))
        stroke = _hex_to_rgb(s["style"].get("strokeColor", "#666"))
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = stroke; shape.line.width = Pt(0.8)
        self._tf_cfg(shape, s["value"], font_size=self._font_pt(8),
                     color=RGBColor(0x36, 0x39, 0x3D))

    def _box(self, slide, s):
        rounded = s["style"].get("rounded") == "1"
        ms = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(
            ms, self._x(s["x"]), self._y(s["y"]),
            self._w(s["w"]), self._h(s["h"]),
        )
        fill = _hex_to_rgb(s["style"].get("fillColor", "#fff"))
        stroke = _hex_to_rgb(s["style"].get("strokeColor", "#666"))
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = stroke; shape.line.width = Pt(0.8)
        fs = int(s["style"].get("fontSize", "10"))
        bold = s["style"].get("fontStyle") == "1"
        self._tf_cfg(shape, s["value"], font_size=self._font_pt(fs),
                     bold=bold, color=RGBColor(0x36, 0x39, 0x3D))

    def _ellipse(self, slide, s):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            self._x(s["x"]), self._y(s["y"]),
            self._w(s["w"]), self._h(s["h"]),
        )
        fill = _hex_to_rgb(s["style"].get("fillColor", "#fff"))
        stroke = _hex_to_rgb(s["style"].get("strokeColor", "#666"))
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = stroke; shape.line.width = Pt(0.8)

    def _edge(self, slide, s, id_geo):
        """绘制多段折线 — 解析 waypoints 构建完整路径，逐段画线，末段加箭头。

        路径构建：
        1. 有 source/target 引用 → 从源/目标形状的 exit/entry 点计算首尾
        2. 有 sourcePoint/targetPoint → 用绝对坐标
        3. waypoints → 中间拐点
        """
        st = s["style"]
        stroke = _hex_to_rgb(st.get("strokeColor", "#666"))
        dashed = st.get("dashed") == "1"
        # 箭头类型
        end_arrow = st.get("endArrow", "classic")
        start_arrow = st.get("startArrow", "none")

        # ── 构建路径点列表 [(x, y), ...] ──────────────────────────────
        pts: list[tuple[float, float]] = []

        # 起点
        src_id = s.get("source")
        src_pt = s.get("source_pt")
        if src_id and src_id in id_geo:
            src = id_geo[src_id]
            ex = float(st.get("exitX", 1.0)); ey = float(st.get("exitY", 0.5))
            pts.append((src["x"] + src["w"] * ex, src["y"] + src["h"] * ey))
        elif src_pt:
            pts.append(src_pt)
        else:
            return  # 无法确定起点

        # 中间拐点
        for wp in s.get("waypoints", []):
            pts.append(wp)

        # 终点
        tgt_id = s.get("target")
        tgt_pt = s.get("target_pt")
        if tgt_id and tgt_id in id_geo:
            tgt = id_geo[tgt_id]
            ix = float(st.get("entryX", 0.0)); iy = float(st.get("entryY", 0.5))
            pts.append((tgt["x"] + tgt["w"] * ix, tgt["y"] + tgt["h"] * iy))
        elif tgt_pt:
            pts.append(tgt_pt)
        else:
            return  # 无法确定终点

        if len(pts) < 2:
            return

        # ── 逐段绘制 ─────────────────────────────────────────────────
        for i in range(len(pts) - 1):
            x1, y1 = pts[i]
            x2, y2 = pts[i + 1]
            is_last = (i == len(pts) - 2)

            seg = slide.shapes.add_connector(
                MSO_CONNECTOR_TYPE.STRAIGHT,
                self._x(x1), self._y(y1), self._x(x2), self._y(y2),
            )
            seg.line.color.rgb = stroke
            seg.line.width = Pt(1.0)

            # 末段加箭头
            if is_last and end_arrow != "none":
                self._add_arrow(seg, "tailEnd")
            # 首段加起始箭头（双向箭头）
            if i == 0 and start_arrow != "none":
                self._add_arrow(seg, "headEnd")

            if dashed:
                spPr = seg._element.find(qn("a:spPr"))
                if spPr is not None:
                    ln = spPr.find(qn("a:ln"))
                    if ln is None:
                        ln = ET.SubElement(spPr, qn("a:ln"))
                    ET.SubElement(ln, qn("a:prstDash")).set("val", "dash")

        # ── 边标签 ──────────────────────────────────────────────────
        if s["value"]:
            mx, my = pts[0]
            txBox = slide.shapes.add_textbox(
                self._x(mx + 4), self._y(my - 14),
                self._w(200), self._h(16),
            )
            tf = txBox.text_frame; tf.word_wrap = False
            p = tf.paragraphs[0]; p.text = s["value"]
            for r in p.runs:
                r.font.size = Pt(8)
                r.font.color.rgb = stroke

    def _add_arrow(self, connector, end: str) -> None:
        spPr = connector._element.find(qn("a:spPr"))
        if spPr is not None:
            ln = spPr.find(qn("a:ln"))
            if ln is None:
                ln = ET.SubElement(spPr, qn("a:ln"))
            el = ET.SubElement(ln, qn(f"a:{end}"))
            el.set("type", "triangle")
            el.set("w", "med")
            el.set("len", "med")

    # ── 入口 ──────────────────────────────────────────────────────────

    def convert(self, input_path: str, output_path: str) -> None:
        shapes, bounds, id_geo = parse_drawio(input_path)
        self._init(bounds)

        cw_in = bounds["w"] / 72.0; ch_in = bounds["h"] / 72.0
        print(f"[INFO] 内容: {bounds['w']:.0f}×{bounds['h']:.0f} du"
              f"  →  {cw_in:.1f}×{ch_in:.1f} in"
              f"  →  缩放 {self.scale:.2f}x  →  "
              f"{cw_in*self.scale:.1f}×{ch_in*self.scale:.1f} in on slide")

        prs = Presentation()
        prs.slide_width = int(SLIDE_W * 914400)
        prs.slide_height = int(SLIDE_H * 914400)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        draw_order = {"bg": 0, "label": 1, "dashed_box": 2,
                      "box": 3, "cylinder": 3, "text": 3}
        ordered = sorted(shapes, key=lambda s: draw_order.get(s["type"], 99))

        counts: dict[str, int] = {}
        draw_fn = {
            "bg": self._bg, "label": self._label,
            "dashed_box": self._dashed_box, "text": self._text,
            "cylinder": self._cylinder, "ellipse": self._ellipse,
            "box": self._box,
        }
        for s in ordered:
            t = s["type"]; counts[t] = counts.get(t, 0) + 1
            fn = draw_fn.get(t)
            if fn:
                fn(slide, s)

        for s in shapes:
            if s["type"] == "edge":
                self._edge(slide, s, id_geo)

        prs.save(output_path)
        print(f"[OK] {output_path}  元素: {counts}")


# ═════════════════════════════════════════════════════════════════════════
# CLI
# ═════════════════════════════════════════════════════════════════════════


def main():
    parser = argparse.ArgumentParser(description="draw.io XML → PowerPoint")
    parser.add_argument("-i", "--input", default="docs/err.xml")
    parser.add_argument("-o", "--output", default=None)
    args = parser.parse_args()

    inp = Path(args.input)
    if not inp.exists():
        fallback = Path("C:/Users/Mr.Wu/Downloads") / inp.name
        if fallback.exists():
            inp = fallback
        else:
            print(f"[FAIL] not found: {args.input}"); return 1
    out = args.output or str(inp.with_suffix(".pptx"))

    Converter().convert(str(inp), out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
